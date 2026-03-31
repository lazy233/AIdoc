import re
from datetime import datetime

from app.core.config import get_settings
from app.models.schemas import GeneratedArtifact, OutlineBlock, ReportGenerationRequest


class PPTGenerator:
  def __init__(self) -> None:
    self.settings = get_settings()

  def generate(
    self,
    request: ReportGenerationRequest,
    outline: list[OutlineBlock],
    template_name: str,
  ) -> GeneratedArtifact:
    self.settings.ensure_directories()
    timestamp = datetime.now()
    file_stem = f'{self._slugify(request.topic)}-{timestamp.strftime("%Y%m%d-%H%M%S")}'

    try:
      from pptx import Presentation
    except ImportError:
      placeholder_path = self.settings.generated_dir / f'{file_stem}.txt'
      placeholder_path.write_text(
        'python-pptx 未安装，当前仅生成了占位说明文件。安装依赖后即可输出真正的 PPT。\n'
        f'主题: {request.topic}\n模板: {template_name}\n输出格式: {request.outputFormat}\n',
        encoding='utf-8',
      )
      return GeneratedArtifact(
        fileName=placeholder_path.name,
        relativePath=f'generated/{placeholder_path.name}',
        slides=len(outline),
        status='dependency_missing',
        outputFormat='txt',
        createdAt=timestamp,
      )

    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = request.topic
    title_slide.placeholders[1].text = f'模板：{template_name}\n生成时间：{timestamp.strftime("%Y-%m-%d %H:%M")}'

    for block in outline:
      slide = presentation.slides.add_slide(presentation.slide_layouts[1])
      slide.shapes.title.text = block.title
      slide.placeholders[1].text = f'{block.description}\n\n页面建议：{block.emphasis}'

    output_path = self.settings.generated_dir / f'{file_stem}.pptx'
    presentation.save(output_path)

    return GeneratedArtifact(
      fileName=output_path.name,
      relativePath=f'generated/{output_path.name}',
      slides=len(outline) + 1,
      status='created',
      outputFormat='pptx',
      createdAt=timestamp,
    )

  def _slugify(self, value: str) -> str:
    cleaned = re.sub(r'[^0-9A-Za-z\u4e00-\u9fa5]+', '-', value).strip('-')
    return cleaned[:40] or 'report'
