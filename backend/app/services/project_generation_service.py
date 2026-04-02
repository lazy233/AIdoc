from __future__ import annotations

import json
import logging
import re
from copy import deepcopy
from datetime import datetime
from pathlib import Path
from typing import Any
from uuid import uuid4

import json5
from langchain_community.chat_models.tongyi import ChatTongyi
from langchain_core.messages import HumanMessage, SystemMessage

from app.core.config import get_settings
from app.models.schemas import (
  GenerationHistoryPayload,
  ProjectChatMessagePayload,
  ProjectGenerationResponse,
  StudentDetail,
  project_detail_to_payload,
)
from app.services.chapter_binding_service import _dashscope_api_key
from app.services.postgres_repository import postgres_repository
from app.services.student_field_catalog import list_student_field_catalog

logger = logging.getLogger(__name__)

_CHAPTER_PLAN_SYSTEM = """你是资深 PPT 文案编辑。本章幻灯片上**所有可见文字**都由你撰写，依据输入中的「章节卡片」与「事实数据 facts」重新梳理，不得照搬罗列原始数据。

输入说明：
- chapter_card：本章标题、说明、补充想法（供你理解意图）
- chapter_facts：从数据库取出的结构化事实，仅供你理解，**不要**在输出里出现字段英文名、路径、JSON、或「标签：内容」式清单
- components：本模板片段中需要处理的文本类占位，每个有唯一 component_key（格式 page_no:component_order）

写作与版式要求（必须遵守）：
1) 全文使用自然、通顺的中文，口吻适合家长会/学业报告类演示。
2) 符合正常 PPT：标题类占位宜短（一般不超过一行）；正文类占位用短句或 2～4 条要点，每条一行，避免长篇论文式段落。
3) **禁止重复**：同一事实或同义表述不要在多个文本框里重复出现；信息只放在最合适的一个框里。
4) **禁止格式化数据**：不要输出「姓名：」「学号：」这类键值对罗列、不要表格感、不要 JSON/Markdown、不要英文字段名或 snake_case。
5) 对 facts 中的数字与名词，融入叙述中表达（例如「总课时 96 学时」），不要逐条对照字段抄录。
6) 你必须为输入中列出的**每一个** component_key 给出一条计划：replace_text（改写占位）、emphasize（需要加粗语气的关键一句）、keep（仅当该占位是装饰性固定文案且必须保留原模板字时）、hide（本页不需要展示该占位时）。
7) replace_text / emphasize 时 text 必填；hide 时 text 为空字符串即可。
8) 输出必须是合法 JSON 对象，禁止 Markdown 代码块。

输出格式：
{
  "component_plans": [
    {"component_key":"3:1", "action":"replace_text", "text":"……"}
  ]
}
"""

_REPAIR_PLAN_SYSTEM = """你是 PPT 文案编辑。上一轮输出缺少部分 component_key。请**仅**为「待补全列表」中的 key 各写一条计划，规则与主任务相同：全中文、PPT 口吻、不重复已有框内语义、禁止键值对罗列与英文字段名。
输出 JSON：{"component_plans":[...]} ，不得包含 Markdown。"""

_FIELD_LABELS = {item['id']: item['label'] for item in list_student_field_catalog()}
_FIELD_IDS_BY_LEN = sorted(_FIELD_LABELS.keys(), key=len, reverse=True)


def _slugify(value: str) -> str:
  cleaned = re.sub(r'[^0-9A-Za-z\u4e00-\u9fa5]+', '-', value).strip('-')
  return cleaned[:40] or 'report'


def _field_value(student: StudentDetail, field_id: str):
  current = student
  for part in field_id.split('.'):
    if isinstance(current, dict):
      current = current.get(part)
    else:
      current = getattr(current, part, None)
    if current is None:
      return None
  return current


def _clean_text(text: str, *, remove_field_prefix: bool = False) -> str:
  cleaned = text.replace('\r', '\n').replace('�', '')
  for field_id in _FIELD_IDS_BY_LEN:
    if field_id in cleaned:
      cleaned = cleaned.replace(field_id, _FIELD_LABELS[field_id])
  cleaned = re.sub(r'[`]{3,}', '', cleaned)
  cleaned = re.sub(r'\b[A-Za-z_][A-Za-z0-9_]*\.[A-Za-z0-9_.]+\b', '', cleaned)
  cleaned = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F]', '', cleaned)
  lines = [line.strip() for line in cleaned.split('\n') if line.strip()]
  if remove_field_prefix:
    lines = [line for line in lines if not re.match(r'^[A-Za-z_][A-Za-z0-9_.]{1,60}\s*[:：]', line)]
  merged = '\n'.join(lines)
  merged = re.sub(r'[ \t]{2,}', ' ', merged).strip()
  return merged


def _contains_chinese(text: str) -> bool:
  return bool(re.search(r'[\u4e00-\u9fff]', text))


def _resolve_section_groups(template):
  pages = sorted(template.pages, key=lambda p: p.page_no)
  if not pages:
    return []

  with_section_order = [p for p in pages if p.section_order is not None]
  if with_section_order:
    grouped: dict[int, list] = {}
    for page in with_section_order:
      grouped.setdefault(int(page.section_order), []).append(page)
    return [grouped[key] for key in sorted(grouped.keys())]

  if template.sections:
    groups = []
    for section in sorted(template.sections, key=lambda s: s.section_order):
      start = section.start_page_no or min(p.page_no for p in pages)
      end = section.end_page_no or max(p.page_no for p in pages)
      section_pages = [p for p in pages if start <= p.page_no <= end]
      if section_pages:
        groups.append(section_pages)
    if groups:
      return groups

  return [[p] for p in pages]


def _chapter_facts_for_llm(project_page, student: StudentDetail) -> list[dict[str, str]]:
  """供大模型理解的结构化事实，不直接写入幻灯片。"""
  facts: list[dict[str, str]] = []
  for binding in sorted(project_page.bindings, key=lambda b: b.field_order):
    value = _field_value(student, binding.field_name)
    if value in (None, '', []):
      continue
    if isinstance(value, list):
      value = '、'.join(str(v) for v in value[:12])
    label = _FIELD_LABELS.get(binding.field_name) or binding.field_name.split('.')[-1]
    raw = _clean_text(str(value), remove_field_prefix=True)[:800]
    if not raw:
      continue
    facts.append({'field_id': binding.field_name, 'label': label, 'value': raw})
  if project_page.manual_text and project_page.manual_text.strip():
    mt = _clean_text(project_page.manual_text.strip(), remove_field_prefix=True)[:800]
    if mt:
      facts.append({'field_id': 'manual', 'label': '编者补充', 'value': mt})
  return facts


def _llm_slide_body_no_components(chapter, chapter_facts: list[dict[str, str]], page_no: int) -> str:
  """模板页未解析出文本组件时，由模型生成一段本章摘要正文。"""
  key = _dashscope_api_key()
  if not key:
    return '本章内容待补充。'
  settings = get_settings()
  llm = ChatTongyi(model=settings.qwen_binding_model, api_key=key, temperature=0.35)
  payload = {
    'chapter_card': {
      'page_order': chapter.page_order,
      'title': chapter.title,
      'description': chapter.description,
      'manual_text': chapter.manual_text,
    },
    'chapter_facts': chapter_facts,
    'template_page_no': page_no,
    'instruction': '本页无单独文本占位解析结果，请写一段适合放入一个正文文本框的中文（3～6 个短句或要点分行），概括本章要点；禁止键值对罗列与重复套话。',
  }
  try:
    result = llm.invoke(
      [
        SystemMessage(
          content='你是 PPT 文案编辑。只输出纯中文正文，禁止 Markdown、禁止 JSON、禁止「标签：值」罗列。'
        ),
        HumanMessage(content=json.dumps(payload, ensure_ascii=False)),
      ]
    )
    text = _coerce_ai_content_to_text(getattr(result, 'content', result))
    text = _clean_text(text, remove_field_prefix=True)[:1200]
    return text or '本章内容待补充。'
  except Exception as exc:  # noqa: BLE001
    logger.warning('无文本组件页 LLM 生成失败: %s', exc)
    return '本章内容待补充。'


def _slice_first_brace_object(s: str, start: int) -> str | None:
  depth = 0
  in_string = False
  escape = False
  quote = ''
  i = start
  while i < len(s):
    c = s[i]
    if in_string:
      if escape:
        escape = False
      elif c == '\\':
        escape = True
      elif c == quote:
        in_string = False
        quote = ''
      i += 1
      continue
    if c in '"\'':
      in_string = True
      quote = c
      i += 1
      continue
    if c == '{':
      depth += 1
    elif c == '}':
      depth -= 1
      if depth == 0:
        return s[start : i + 1]
    i += 1
  return None


def _coerce_ai_content_to_text(content: Any) -> str:
  if content is None:
    return ''
  if isinstance(content, str):
    return content
  if isinstance(content, list):
    chunks: list[str] = []
    for part in content:
      if isinstance(part, str):
        chunks.append(part)
      elif isinstance(part, dict):
        t = part.get('text')
        if isinstance(t, str):
          chunks.append(t)
    return ''.join(chunks) if chunks else json.dumps(content, ensure_ascii=False)
  return str(content)


def _parse_json_object(text: str) -> dict[str, Any]:
  raw = text.strip()
  if raw.startswith('```'):
    raw = raw[3:].lstrip()
    if raw[:4].lower() == 'json':
      raw = raw[4:].lstrip()
    raw = raw.rstrip()
    if raw.endswith('```'):
      raw = raw[:-3].rstrip()
  start = raw.find('{')
  if start == -1:
    raise ValueError('模型未返回 JSON 对象')
  chunk = _slice_first_brace_object(raw, start)
  if not chunk:
    raise ValueError('模型 JSON 括号未闭合')
  obj = json5.loads(chunk)
  if not isinstance(obj, dict):
    raise ValueError('模型 JSON 根节点需为对象')
  return obj


def _rows_to_plan(rows: list[Any], allowed: set[str]) -> dict[str, dict[str, str]]:
  allowed_actions = {'replace_text', 'keep', 'hide', 'emphasize'}
  plan: dict[str, dict[str, str]] = {}
  if not isinstance(rows, list):
    return plan
  for row in rows:
    if not isinstance(row, dict):
      continue
    comp_key = row.get('component_key')
    action = row.get('action', 'replace_text')
    text = row.get('text', '')
    if not (isinstance(comp_key, str) and comp_key in allowed):
      continue
    if not isinstance(action, str) or action not in allowed_actions:
      action = 'replace_text'
    if not isinstance(text, str):
      text = ''
    clean_text = _clean_text(text, remove_field_prefix=True)[:800]
    if action == 'hide':
      plan[comp_key] = {'action': 'hide', 'text': ''}
      continue
    if action in {'replace_text', 'emphasize'} and not clean_text:
      continue
    if action == 'keep':
      plan[comp_key] = {'action': 'keep', 'text': clean_text}
      continue
    plan[comp_key] = {'action': action, 'text': clean_text}
  return plan


def _repair_llm_plan(
  llm: ChatTongyi,
  chapter,
  components: list[dict[str, Any]],
  chapter_facts: list[dict[str, str]],
  existing: dict[str, dict[str, str]],
  missing: set[str],
) -> dict[str, dict[str, str]]:
  comp_subset = [c for c in components if c['component_key'] in missing]
  summary = {
    k: {'action': v['action'], 'text_preview': (v.get('text') or '')[:120]}
    for k, v in existing.items()
    if k not in missing
  }
  payload = {
    'chapter_card': {
      'page_order': chapter.page_order,
      'title': chapter.title,
      'description': chapter.description,
      'manual_text': chapter.manual_text,
    },
    'chapter_facts': chapter_facts,
    'already_filled': summary,
    'missing_component_keys': sorted(missing),
    'components_to_fill': comp_subset,
  }
  result = llm.invoke(
    [
      SystemMessage(content=_REPAIR_PLAN_SYSTEM),
      HumanMessage(content=json.dumps(payload, ensure_ascii=False)),
    ]
  )
  content = _coerce_ai_content_to_text(getattr(result, 'content', result))
  obj = _parse_json_object(content)
  rows = obj.get('component_plans')
  if not isinstance(rows, list):
    rows = obj.get('component_texts')
  return _rows_to_plan(rows if isinstance(rows, list) else [], missing)


def _build_llm_plan_for_chapter(
  chapter, section_pages, chapter_facts: list[dict[str, str]]
) -> dict[str, dict[str, str]]:
  key = _dashscope_api_key()
  if not key:
    raise ValueError('未配置大模型 API Key，无法按文案生成 PPT')

  settings = get_settings()
  llm = ChatTongyi(model=settings.qwen_binding_model, api_key=key, temperature=0.35)

  components: list[dict[str, Any]] = []
  for page in section_pages:
    for comp in sorted(page.components, key=lambda c: c.component_order):
      # 即使模板占位符本身没有默认文字，也要纳入填充范围；
      # 否则会出现生成时“未解析到文本组件”，进而走回退逻辑。
      if not (
        comp.placeholder_key is not None
        or comp.text_content
        or (comp.component_name and '文本' in comp.component_name)
      ):
        continue
      components.append(
        {
          'component_key': f'{page.page_no}:{comp.component_order}',
          'page_no': page.page_no,
          'component_order': comp.component_order,
          'component_name': comp.component_name,
          'placeholder_key': comp.placeholder_key,
          'binding_hint': comp.binding_hint,
          'current_text': (comp.text_content or '')[:300],
        }
      )

  if not components:
    return {}

  allowed = {item['component_key'] for item in components}
  payload = {
    'chapter_card': {
      'page_order': chapter.page_order,
      'title': chapter.title,
      'description': chapter.description,
      'manual_text': chapter.manual_text,
    },
    'chapter_facts': chapter_facts,
    'components': components,
  }

  plan: dict[str, dict[str, str]] = {}
  last_error: Exception | None = None
  for _attempt in range(3):
    try:
      result = llm.invoke(
        [
          SystemMessage(content=_CHAPTER_PLAN_SYSTEM),
          HumanMessage(content=json.dumps(payload, ensure_ascii=False)),
        ]
      )
      content = _coerce_ai_content_to_text(getattr(result, 'content', result))
      obj = _parse_json_object(content)
      rows = obj.get('component_plans')
      if not isinstance(rows, list):
        rows = obj.get('component_texts')
      if not isinstance(rows, list):
        last_error = ValueError('缺少 component_plans 数组')
        continue
      plan = _rows_to_plan(rows, allowed)
      if len(plan) >= len(allowed):
        break
      last_error = ValueError(f'component 覆盖不全: {len(plan)}/{len(allowed)}')
    except Exception as exc:  # noqa: BLE001
      last_error = exc
      continue

  missing = allowed - set(plan.keys())
  for _ in range(2):
    if not missing:
      break
    try:
      patch = _repair_llm_plan(llm, chapter, components, chapter_facts, plan, missing)
      plan.update(patch)
      missing = allowed - set(plan.keys())
    except Exception as exc:  # noqa: BLE001
      logger.warning('章节 %s 补全计划失败: %s', chapter.page_order, exc)
      break

  if missing:
    for mk in missing:
      meta = next((c for c in components if c['component_key'] == mk), None)
      meta = meta or {}

      binding_hint = meta.get('binding_hint')
      title = getattr(chapter, 'title', '') or ''
      facts_values = [f.get('value') for f in chapter_facts if isinstance(f, dict) and f.get('value')]

      if binding_hint == '标题占位':
        plan[mk] = {'action': 'replace_text', 'text': (title or meta.get('current_text') or '本章标题')[:120]}
      elif binding_hint == '正文占位':
        points = facts_values[:4]
        plan[mk] = {
          'action': 'replace_text',
          'text': ('\n'.join(points) if points else '本章要点已自动补全')[:600],
        }
      elif binding_hint == '副标题占位':
        desc = getattr(chapter, 'description', '') or ''
        plan[mk] = {'action': 'replace_text', 'text': (desc or title or meta.get('current_text') or '补充说明')[:160]}
      elif binding_hint == '页眉页脚占位':
        plan[mk] = {'action': 'replace_text', 'text': (title or meta.get('current_text') or '')[:120]}
      else:
        first = facts_values[0] if facts_values else (meta.get('current_text') or '')
        plan[mk] = {'action': 'replace_text', 'text': (first or title or '本页内容已自动补全')[:200]}

  return plan


def _clone_slide(prs, source_slide):
  target = prs.slides.add_slide(source_slide.slide_layout)
  target_sp_tree = target.shapes._spTree
  for shape in list(target.shapes):
    target_sp_tree.remove(shape.element)

  for shape in source_slide.shapes:
    new_el = deepcopy(shape.element)
    target_sp_tree.insert_element_before(new_el, 'p:extLst')

  source_bg = source_slide.element.cSld.bg
  if source_bg is not None:
    target_c_sld = target.element.cSld
    target_bg = target_c_sld.bg
    if target_bg is not None:
      target_c_sld.remove(target_bg)
    target_c_sld.insert(0, deepcopy(source_bg))

  rel_id_map: dict[str, str] = {}
  for rel in source_slide.part.rels.values():
    if 'notesSlide' in rel.reltype or 'slideLayout' in rel.reltype:
      continue
    new_rel_id = target.part.rels._add_relationship(rel.reltype, rel._target, rel.is_external)
    rel_id_map[rel.rId] = new_rel_id

  if rel_id_map:
    for el in target.element.iter():
      for attr_name, attr_value in list(el.attrib.items()):
        if isinstance(attr_value, str) and attr_value in rel_id_map:
          el.set(attr_name, rel_id_map[attr_value])
  return target


def _remove_first_n_slides(prs, n: int) -> None:
  if n <= 0:
    return
  sld_id_lst = prs.slides._sldIdLst
  for _ in range(min(n, len(sld_id_lst))):
    sld_id = sld_id_lst[0]
    rel_id = sld_id.rId
    prs.part.drop_rel(rel_id)
    sld_id_lst.remove(sld_id)


class ProjectGenerationService:
  def __init__(self) -> None:
    self.settings = get_settings()

  def generate(self, project_id: str) -> ProjectGenerationResponse:
    project = postgres_repository.get_project(project_id)
    if not project:
      raise ValueError('项目不存在')
    if not project.ppt_template_id:
      raise ValueError('未选择 PPT 模板')
    if not project.student_id:
      raise ValueError('未选择学生，无法提取字段数据')

    template = postgres_repository.get_ppt_template(project.ppt_template_id)
    if not template:
      raise ValueError('PPT 模板不存在')

    student = postgres_repository.get_student(project.student_id)
    if not student:
      raise ValueError('学生数据不存在')

    section_groups = _resolve_section_groups(template)
    if not section_groups:
      raise ValueError('模板缺少可用页面结构')

    try:
      from pptx import Presentation
    except ImportError as exc:
      raise ValueError('python-pptx 未安装，无法生成 PPT') from exc

    self.settings.ensure_directories()
    source_file = (template.source_file_path or '').strip()
    if not source_file:
      raise ValueError('模板缺少源文件路径，无法按模板生成')
    template_path = Path(source_file)
    if not template_path.is_absolute():
      template_path = self.settings.storage_root / template_path
    if not template_path.exists():
      raise ValueError(f'模板源文件不存在: {template_path}')
    if not _dashscope_api_key():
      raise ValueError('未配置大模型 API Key，无法生成 PPT（全部文案由模型生成）')
    output = Presentation(str(template_path))
    source_slides = list(output.slides)
    if not source_slides:
      raise ValueError('模板源文件没有可用页面')
    original_slide_count = len(source_slides)

    project_pages = sorted(project.pages, key=lambda p: p.page_order)
    chapter_logs: list[str] = []
    chapter_log_rows: list[dict[str, Any]] = []
    for idx, chapter in enumerate(project_pages):
      section_index = min(idx, len(section_groups) - 1)
      section_pages = section_groups[section_index]
      chapter_facts = _chapter_facts_for_llm(chapter, student)

      def _default_fill_text(comp: Any) -> str:
        binding_hint = getattr(comp, 'binding_hint', None)
        title = getattr(chapter, 'title', '') or ''
        facts_values = [f.get('value') for f in chapter_facts if isinstance(f, dict) and f.get('value')]
        if binding_hint == '标题占位':
          return (title or getattr(comp, 'text_content', None) or '本章标题')[:120]
        if binding_hint == '正文占位':
          points = facts_values[:4]
          return ('\n'.join(points) if points else '本章要点已自动补全')[:600]
        if binding_hint == '副标题占位':
          desc = getattr(chapter, 'description', '') or ''
          return (desc or title or getattr(comp, 'text_content', None) or '补充说明')[:160]
        if binding_hint == '页眉页脚占位':
          return (title or getattr(comp, 'text_content', None) or '')[:120]
        first = facts_values[0] if facts_values else (getattr(comp, 'text_content', None) or '')
        return (first or title or '本页内容已自动补全')[:200]

      llm_plan = _build_llm_plan_for_chapter(chapter, section_pages, chapter_facts)
      applied_by_llm = 0
      fallback_count = 0
      for page_idx, template_page in enumerate(section_pages, start=1):
        source_index = max(0, min(template_page.page_no - 1, original_slide_count - 1))
        slide = _clone_slide(output, source_slides[source_index])

        text_components = [
          c
          for c in template_page.components
          if c.placeholder_key is not None or c.text_content or (c.component_name and '文本' in c.component_name)
        ]
        text_components = sorted(text_components, key=lambda c: c.component_order)
        if not text_components:
          body = slide.shapes.add_textbox(685800, 2286000, 10668000, 3657600)
          body.text_frame.text = _llm_slide_body_no_components(chapter, chapter_facts, template_page.page_no)
          applied_by_llm += 1
          continue

        for comp in text_components:
          comp_key = f'{template_page.page_no}:{comp.component_order}'
          shape_idx = comp.component_order - 1
          if shape_idx < 0 or shape_idx >= len(slide.shapes):
            fallback_count += 1
            continue
          shape = slide.shapes[shape_idx]
          if not getattr(shape, 'has_text_frame', False):
            fallback_count += 1
            continue
          slot = llm_plan.get(comp_key)
          if not slot:
            # 不再写入“—”回退符号；尽量用章节 facts 做一个可展示的默认内容。
            value = _default_fill_text(comp)
            value = _clean_text(str(value), remove_field_prefix=True)
            if not value:
              continue
            shape.text_frame.clear()
            shape.text_frame.text = value
            applied_by_llm += 1
            continue

          action = slot.get('action', 'replace_text')
          if action == 'hide':
            shape.element.getparent().remove(shape.element)
            applied_by_llm += 1
            continue

          if action == 'keep':
            # 模板如果本来就没给默认文案，就保持原样，避免清空成“空文本”。
            value = comp.text_content or ''
            if not value:
              continue
          elif action == 'emphasize':
            base = slot.get('text') or ''
            value = f'【重点】{base}' if base else (comp.text_content or '')
          else:
            value = slot.get('text') or ''

          value = _clean_text(str(value), remove_field_prefix=True)
          if not value:
            # 不写入回退符号；让模板原内容（占位符）保持可见。
            continue

          if value and not _contains_chinese(value) and action != 'keep':
            # 非中文时优先保留模板现有文本（若有）。
            fallback_value = comp.text_content or ''
            fallback_value = _clean_text(str(fallback_value), remove_field_prefix=True)
            if fallback_value:
              value = fallback_value
            else:
              continue

          shape.text_frame.clear()
          shape.text_frame.text = value
          applied_by_llm += 1
      chapter_logs.append(
        f'章节{chapter.page_order}: 模板片段{len(section_pages)}，LLM命中{applied_by_llm}，回退{fallback_count}'
      )
      chapter_log_rows.append(
        {
          'chapter_order': chapter.page_order,
          'template_section_index': section_index + 1,
          'template_page_count': len(section_pages),
          'llm_hit_count': applied_by_llm,
          'fallback_count': fallback_count,
          'status': 'completed',
          'error_message': None,
        }
      )

    _remove_first_n_slides(output, original_slide_count)

    timestamp = datetime.now()
    file_name = f"{_slugify(project.prompt or project.report_type)}-{timestamp.strftime('%Y%m%d-%H%M%S')}.pptx"
    output_path = self.settings.generated_dir / file_name
    output.save(output_path)

    history = postgres_repository.create_history_entry(
      GenerationHistoryPayload(
        project_id=project.id,
        student_id=project.student_id,
        report_title=f"{student.name}-{project.report_type}",
        ppt_template_id=project.ppt_template_id,
        report_template_id=project.report_template_id,
        output_format='pptx',
        status='completed',
        output_file_path=f'generated/{file_name}',
      )
    )
    postgres_repository.create_generation_chapter_logs(history.id, chapter_log_rows)

    updated_payload = project_detail_to_payload(project).model_copy(
      update={
        'status': 'completed',
        'messages': [
          *project.messages,
          ProjectChatMessagePayload(
            role='assistant',
            content=f"已完成章节化生成，共 {len(output.slides)} 页，文件：{file_name}。{'；'.join(chapter_logs[:6])}",
          ),
        ],
      }
    )
    postgres_repository.update_project(project.id, updated_payload)

    return ProjectGenerationResponse(
      task_id=f'TASK-{uuid4().hex[:12]}',
      project_id=project.id,
      history_id=history.id,
      status='completed',
      output_file_path=f'generated/{file_name}',
      output_file_name=file_name,
      slide_count=len(output.slides),
      chapter_logs=chapter_logs,
      chapter_log_items=chapter_log_rows,
      message='已按章节生成并合成完整 PPT',
    )


project_generation_service = ProjectGenerationService()
