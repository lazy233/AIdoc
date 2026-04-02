from __future__ import annotations

import json
import logging
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Any
from uuid import uuid4
from zipfile import ZipFile

import json5
from langchain_community.chat_models.tongyi import ChatTongyi
from langchain_core.messages import HumanMessage, SystemMessage

from app.core.config import get_settings
from app.models.schemas import (
  PptTemplateComponentPayload,
  PptTemplateDetail,
  PptTemplatePagePayload,
  PptTemplatePayload,
  PptTemplateSectionPayload,
)
from app.services.postgres_repository import postgres_repository

logger = logging.getLogger(__name__)


def _normalize_template_name(filename: str | None, fallback_stem: str) -> str:
  stem = Path(filename or fallback_stem).stem.strip()
  return stem or fallback_stem


def _guess_aspect_ratio(width: int | None, height: int | None) -> str | None:
  if not width or not height:
    return None

  ratio = width / height
  presets = {
    '16:9': 16 / 9,
    '4:3': 4 / 3,
    '16:10': 16 / 10,
  }
  best = min(presets.items(), key=lambda item: abs(item[1] - ratio))
  if abs(best[1] - ratio) <= 0.08:
    return best[0]
  return f'{width}:{height}'


def _extract_text(shape) -> str | None:
  text = getattr(shape, 'text', None)
  if not text:
    return None
  normalized = ' '.join(str(text).split())
  return normalized[:500] if normalized else None


def _paragraph_lines_from_text_frame(text_frame) -> list[str]:
  if text_frame is None:
    return []
  lines: list[str] = []
  try:
    for para in text_frame.paragraphs:
      parts = []
      for run in para.runs:
        if run.text:
          parts.append(run.text)
      raw = ''.join(parts).strip() if parts else (para.text or '').strip()
      if raw:
        lines.append(' '.join(raw.split()))
  except Exception:
    pass
  return lines


def _collect_shape_text_parts(shape) -> list[str]:
  """递归抽取形状内可见文本（组合内子形状、表格单元格），不改变顶层 shapes 顺序。"""
  try:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
  except ImportError:
    MSO_SHAPE_TYPE = None

  parts: list[str] = []
  st = getattr(shape, 'shape_type', None)
  if MSO_SHAPE_TYPE is not None and st == MSO_SHAPE_TYPE.GROUP:
    try:
      for child in shape.shapes:
        parts.extend(_collect_shape_text_parts(child))
    except Exception:
      pass
    return parts

  if getattr(shape, 'has_table', False):
    try:
      for row in shape.table.rows:
        for cell in row.cells:
          parts.extend(_paragraph_lines_from_text_frame(cell.text_frame))
    except Exception:
      pass
    return parts

  if getattr(shape, 'has_text_frame', False):
    parts.extend(_paragraph_lines_from_text_frame(shape.text_frame))

  return parts


def _deep_extract_text(shape, max_len: int = 800) -> str | None:
  """比 shape.text 更完整：含组合内、表格单元格、多段落的拼接。"""
  chunks = _collect_shape_text_parts(shape)
  if not chunks:
    return _extract_text(shape)
  merged = ' '.join(chunks)
  normalized = ' '.join(merged.split())
  return normalized[:max_len] if normalized else None


def _placeholder_type_name(shape) -> str | None:
  if not getattr(shape, 'is_placeholder', False):
    return None
  try:
    return shape.placeholder_format.type.name
  except Exception:
    return None


def _binding_hint_for_shape(shape) -> str | None:
  """用于 UI/后续按角色匹配；不改变 component_order。"""
  name = _placeholder_type_name(shape)
  if not name:
    return None
  title_like = {'TITLE', 'CENTER_TITLE', 'VERTICAL_TITLE', 'CENTERED_TITLE'}
  if name in title_like:
    return '标题占位'
  if name == 'SUBTITLE':
    return '副标题占位'
  if name in {'BODY', 'OBJECT', 'VERTICAL_BODY', 'CONTENT'} or name.startswith('BODY'):
    return '正文占位'
  if name in {'FOOTER', 'SLIDE_NUMBER', 'DATE', 'HEADER'}:
    return '页眉页脚占位'
  return '其它占位'


def _slide_text_sort_priority(shape) -> int:
  hint = _binding_hint_for_shape(shape)
  if hint == '标题占位':
    return 0
  if hint == '副标题占位':
    return 1
  if hint == '正文占位':
    return 2
  return 3


_TOC_HINTS = ('目录', '目 录', 'contents', 'table of contents', 'agenda', 'overview', 'index')
_SECTION_HINTS = ('章节', '第', 'chapter', 'part', 'section', '模块', '单元')
_COVER_HINTS = ('封面', 'title', '首页')


def _extract_heading_candidate(slide_texts: list[str]) -> str | None:
  for text in slide_texts:
    normalized = text.strip()
    if not normalized:
      continue
    if len(normalized) <= 60:
      return normalized
  return slide_texts[0].strip()[:60] if slide_texts else None


def _is_toc_page(layout_name: str | None, joined_text: str, text_count: int) -> bool:
  layout = (layout_name or '').lower()
  if any(hint in layout for hint in ('toc', 'content', 'agenda', '目录')):
    return True
  if any(hint in joined_text for hint in _TOC_HINTS):
    return True
  # 目录页通常有较多短文本条目
  return text_count >= 5 and ('第' in joined_text and ('章' in joined_text or '节' in joined_text))


def _is_section_page(layout_name: str | None, heading: str | None, joined_text: str, component_count: int) -> bool:
  if not heading:
    return False
  layout = (layout_name or '').lower()
  heading_lower = heading.lower()
  if any(hint in layout for hint in ('section', 'chapter', 'divider', '章节')):
    return True
  if any(hint in heading_lower for hint in _COVER_HINTS):
    return False
  if re.search(r'^第[一二三四五六七八九十0-9]+[章节部分篇]', heading):
    return True
  if re.search(r'^(chapter|part|section)\s*[0-9ivx]+', heading_lower):
    return True
  if any(hint in heading_lower for hint in _SECTION_HINTS) and component_count <= 8:
    return True
  return any(hint in joined_text for hint in ('本章', '本节', '学习目标', '章节目标')) and component_count <= 10


def _classify_page_type(layout_name: str | None, slide_texts: list[str], component_count: int) -> tuple[str, str | None]:
  joined_text = ' '.join(slide_texts).lower()
  heading = _extract_heading_candidate(slide_texts)
  if _is_toc_page(layout_name, joined_text, len(slide_texts)):
    return 'toc', heading
  if _is_section_page(layout_name, heading, joined_text, component_count):
    return 'section', heading
  return 'content', heading


def _read_windows_env_registry(*names: str) -> str | None:
  if os.name != 'nt':
    return None
  try:
    import winreg
  except ImportError:
    return None

  locations = (
    (winreg.HKEY_CURRENT_USER, r'Environment'),
    (winreg.HKEY_LOCAL_MACHINE, r'SYSTEM\CurrentControlSet\Control\Session Manager\Environment'),
  )
  wanted = {name.lower() for name in names}
  for root, path in locations:
    try:
      with winreg.OpenKey(root, path) as key:
        index = 0
        while True:
          try:
            name, value, _value_type = winreg.EnumValue(key, index)
          except OSError:
            break
          if name.lower() in wanted and str(value).strip():
            return str(value).strip()
          index += 1
    except OSError:
      continue
  return None


def _dashscope_api_key() -> str | None:
  settings = get_settings()
  for candidate in (
    settings.dashscope_api_key,
    os.environ.get('DASHSCOPE_API_KEY'),
    os.environ.get('OPENAI_API_KEY'),
    os.environ.get('OPENAIAPIKEY'),
    os.environ.get('openaiapikey'),
    _read_windows_env_registry('DASHSCOPE_API_KEY', 'OPENAI_API_KEY', 'OPENAIAPIKEY', 'openaiapikey'),
  ):
    if candidate and str(candidate).strip():
      return str(candidate).strip()
  return None


def _slice_first_brace_object(s: str, start: int) -> str | None:
  depth = 0
  in_string = False
  escape = False
  quote = ''
  i = start
  while i < len(s):
    c = s[i]
    if in_string:
      if escape:
        escape = False
      elif c == '\\':
        escape = True
      elif c == quote:
        in_string = False
        quote = ''
      i += 1
      continue
    if c in '"\'':
      in_string = True
      quote = c
      i += 1
      continue
    if c == '{':
      depth += 1
    elif c == '}':
      depth -= 1
      if depth == 0:
        return s[start : i + 1]
    i += 1
  return None


def _parse_json_object(text: str) -> dict[str, Any]:
  raw = text.strip()
  if raw.startswith('```'):
    raw = raw[3:].lstrip()
    if raw[:4].lower() == 'json':
      raw = raw[4:].lstrip()
    raw = raw.lstrip('\n\r')
    if raw.rstrip().endswith('```'):
      raw = raw.rstrip()[:-3].rstrip()
  start = raw.find('{')
  if start < 0:
    raise ValueError('模型未返回 JSON 对象')
  chunk = _slice_first_brace_object(raw, start)
  if not chunk:
    raise ValueError('模型返回 JSON 不完整')
  obj = json5.loads(chunk)
  if not isinstance(obj, dict):
    raise ValueError('模型返回 JSON 根节点不是对象')
  return obj


_PPT_PAGE_CLASSIFY_SYSTEM = """你是 PPT 结构识别助手。
给定每页的：版式名 layout_name、占位符类型摘要 placeholder_roles（如 TITLE,BODY）、标题候选、文本预览、组件数量、规则初判 rule_page_type。请输出每页类型：
- toc: 目录页
- section: 章节扉页/章节标题页（通常用于引出后续几页）
- content: 正文页

要求：
1) 仅返回 JSON 对象，不要 Markdown。
2) page_type 只能是 toc/section/content。
3) 若某页为 section，section_title 填章节名；否则 section_title 置空字符串。
4) 判定章节扉页时优先参考：占位以 TITLE/CENTER_TITLE 为主且正文较少、或版式名含 section/chapter/divider/章节。

输出格式严格如下：
{"pages":[{"page_no":1,"page_type":"content","section_title":""}]}
"""


def _llm_refine_page_types(page_candidates: list[dict[str, Any]]) -> dict[int, tuple[str, str | None]]:
  key = _dashscope_api_key()
  if not key:
    return {}
  settings = get_settings()
  llm = ChatTongyi(
    model=settings.qwen_binding_model,
    api_key=key,
    temperature=0.1,
  )
  payload = {
    'instruction': '请做页面类型判定，优先识别章节扉页',
    'pages': page_candidates,
  }
  try:
    result = llm.invoke([SystemMessage(content=_PPT_PAGE_CLASSIFY_SYSTEM), HumanMessage(content=json.dumps(payload, ensure_ascii=False))])
    content = getattr(result, 'content', result)
    text = content if isinstance(content, str) else str(content)
    obj = _parse_json_object(text)
    rows = obj.get('pages')
    if not isinstance(rows, list):
      return {}
    refined: dict[int, tuple[str, str | None]] = {}
    for row in rows:
      if not isinstance(row, dict):
        continue
      page_no = row.get('page_no')
      page_type = row.get('page_type')
      if not isinstance(page_no, int):
        continue
      if page_type not in ('toc', 'section', 'content'):
        continue
      section_title = row.get('section_title')
      refined[page_no] = (page_type, section_title.strip() if isinstance(section_title, str) and section_title.strip() else None)
    return refined
  except Exception as exc:
    logger.warning('PPT 页面大模型判定失败，回退规则判定: %s', exc)
    return {}


def _extract_embedded_cover_image(source_path: Path) -> str | None:
  settings = get_settings()
  preview_dir = settings.templates_dir / 'covers'
  preview_dir.mkdir(parents=True, exist_ok=True)

  try:
    with ZipFile(source_path) as archive:
      for candidate in ('docProps/thumbnail.jpeg', 'docProps/thumbnail.jpg', 'docProps/thumbnail.png'):
        if candidate in archive.namelist():
          suffix = Path(candidate).suffix or '.jpg'
          target_path = preview_dir / f'{uuid4().hex}{suffix}'
          target_path.write_bytes(archive.read(candidate))
          return target_path.relative_to(settings.storage_root).as_posix()
  except Exception:
    return None

  return None


def _parse_presentation_metadata(source_path: Path) -> tuple[PptTemplatePayload, str | None]:
  from pptx import Presentation

  parsed_at = datetime.utcnow()
  presentation = Presentation(str(source_path))
  core = presentation.core_properties
  slides = list(presentation.slides)

  pages: list[PptTemplatePagePayload] = []
  sections: list[PptTemplateSectionPayload] = []
  pages_json: list[dict[str, object]] = []
  components_json: list[dict[str, object]] = []
  draft_pages: list[dict[str, Any]] = []

  for page_no, slide in enumerate(slides, start=1):
    slide_components: list[PptTemplateComponentPayload] = []
    slide_text_entries: list[tuple[int, int, str]] = []
    placeholder_type_names: list[str] = []

    for component_order, shape in enumerate(slide.shapes, start=1):
      component_type = getattr(getattr(shape, 'shape_type', None), 'name', 'shape').lower()
      placeholder_idx = None
      ph_name = _placeholder_type_name(shape)
      if ph_name:
        placeholder_type_names.append(ph_name)
      if getattr(shape, 'is_placeholder', False):
        try:
          placeholder_idx = str(shape.placeholder_format.idx)
        except Exception:
          placeholder_idx = 'placeholder'
      text_content = _deep_extract_text(shape)
      binding_hint = _binding_hint_for_shape(shape)
      if text_content:
        slide_text_entries.append((_slide_text_sort_priority(shape), component_order, text_content))

      line_count = len(_collect_shape_text_parts(shape))
      component = PptTemplateComponentPayload(
        component_order=component_order,
        component_type=component_type,
        component_name=getattr(shape, 'name', None),
        placeholder_key=placeholder_idx,
        text_content=text_content,
        x=float(getattr(shape, 'left', 0) or 0),
        y=float(getattr(shape, 'top', 0) or 0),
        width=float(getattr(shape, 'width', 0) or 0),
        height=float(getattr(shape, 'height', 0) or 0),
        binding_hint=binding_hint,
        raw_component_json={
          'shape_type': component_type,
          'name': getattr(shape, 'name', None),
          'has_text': bool(text_content),
          'placeholder_type': ph_name,
          'text_line_count': line_count,
          'parse': 'deep_v2',
        },
      )
      slide_components.append(component)
      components_json.append(
        {
          'page_no': page_no,
          'component_order': component_order,
          'component_type': component_type,
          'component_name': getattr(shape, 'name', None),
          'placeholder_key': placeholder_idx,
          'placeholder_type': ph_name,
          'binding_hint': binding_hint,
          'text_excerpt': text_content,
        }
      )

    slide_text_entries.sort(key=lambda item: (item[0], item[1]))
    slide_texts = [entry[2] for entry in slide_text_entries]
    placeholder_summary = ','.join(dict.fromkeys(placeholder_type_names))

    layout_name = getattr(getattr(slide, 'slide_layout', None), 'name', None)
    rule_page_type, heading = _classify_page_type(layout_name, slide_texts, len(slide_components))
    draft_pages.append(
      {
        'page_no': page_no,
        'layout_name': layout_name,
        'component_count': len(slide_components),
        'slide_texts': slide_texts,
        'heading': heading,
        'rule_page_type': rule_page_type,
        'components': slide_components,
        'placeholder_summary': placeholder_summary,
      }
    )

  llm_input_pages = [
    {
      'page_no': row['page_no'],
      'layout_name': row['layout_name'],
      'component_count': row['component_count'],
      'heading_candidate': row['heading'] or '',
      'text_preview': ' | '.join(row['slide_texts'][:8]),
      'rule_page_type': row['rule_page_type'],
      'placeholder_roles': row.get('placeholder_summary') or '',
    }
    for row in draft_pages
  ]
  llm_refined = _llm_refine_page_types(llm_input_pages)
  current_section_order: int | None = None

  for row in draft_pages:
    page_no = row['page_no']
    layout_name = row['layout_name']
    component_count = row['component_count']
    heading = row['heading']
    slide_texts = row['slide_texts']
    page_type, section_title = llm_refined.get(page_no, (row['rule_page_type'], None))
    if page_type == 'section':
      current_section_order = len(sections) + 1
      sections.append(
        PptTemplateSectionPayload(
          section_order=current_section_order,
          section_title=section_title or heading or f'章节 {current_section_order}',
          start_page_no=page_no,
          end_page_no=page_no,
          page_count=1,
          summary=None,
        )
      )
    elif current_section_order is not None:
      section = sections[current_section_order - 1]
      section.end_page_no = page_no
      section.page_count = (section.end_page_no or page_no) - (section.start_page_no or page_no) + 1

    page = PptTemplatePagePayload(
      section_order=current_section_order,
      page_no=page_no,
      page_title=heading or f'第 {page_no} 页',
      page_type=page_type,
      layout_name=layout_name,
      component_count=component_count,
      raw_page_json={
        'page_no': page_no,
        'section_order': current_section_order,
        'page_type': page_type,
        'layout_name': layout_name,
        'component_count': component_count,
        'text_preview': (slide_texts[0] if slide_texts else None),
        'placeholder_summary': row.get('placeholder_summary') or '',
        'rule_page_type': row['rule_page_type'],
        'llm_refined': page_no in llm_refined,
        'parse': 'deep_v2',
      },
      components=row['components'],
    )
    pages.append(page)
    pages_json.append(
      {
        'page_no': page_no,
        'section_order': current_section_order,
        'page_type': page_type,
        'layout_name': layout_name,
        'component_count': component_count,
      }
    )

  theme_name = (core.title or core.subject or '').strip() or None
  payload = PptTemplatePayload(
    name=_normalize_template_name(source_path.name, source_path.stem),
    category=None,
    status='draft',
    source_file_name=source_path.name,
    source_file_path=source_path.relative_to(get_settings().storage_root).as_posix(),
    cover_image_path=_extract_embedded_cover_image(source_path),
    file_size=source_path.stat().st_size,
    page_count=len(slides),
    aspect_ratio=_guess_aspect_ratio(getattr(presentation, 'slide_width', None), getattr(presentation, 'slide_height', None)),
    theme_name=theme_name,
    template_version='v2',
    parse_status='parsed',
    parse_error=None,
    outline_json=[],
    pages_json=pages_json,
    components_json=components_json,
    style_tokens_json={},
    slot_bindings_json={},
    parsed_at=parsed_at,
    sections=sections,
    pages=pages,
  )
  return payload, None


def ingest_ppt_template(file_name: str, content: bytes) -> PptTemplateDetail:
  settings = get_settings()
  settings.ensure_directories()

  suffix = Path(file_name).suffix or '.pptx'
  stored_path = settings.templates_dir / f'{uuid4().hex}{suffix}'
  stored_path.write_bytes(content)

  try:
    payload, _warning = _parse_presentation_metadata(stored_path)
  except Exception as exc:
    payload = PptTemplatePayload(
      name=_normalize_template_name(file_name, stored_path.stem),
      category=None,
      status='draft',
      source_file_name=file_name,
      source_file_path=stored_path.relative_to(settings.storage_root).as_posix(),
      cover_image_path=None,
      file_size=len(content),
      page_count=0,
      aspect_ratio=None,
      theme_name=None,
      template_version='v2',
      parse_status='failed',
      parse_error=str(exc),
      outline_json=[],
      pages_json=[],
      components_json=[],
      style_tokens_json={},
      slot_bindings_json={},
      parsed_at=datetime.utcnow(),
      sections=[],
      pages=[],
    )

  payload = payload.model_copy(
    update={
      'source_file_name': file_name,
      'source_file_path': stored_path.relative_to(settings.storage_root).as_posix(),
      'file_size': len(content),
      'name': _normalize_template_name(file_name, stored_path.stem),
    }
  )
  return postgres_repository.create_ppt_template(payload)
